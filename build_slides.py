# -*- coding: utf-8 -*-
"""
AI 에이전트와 응용 — 5일 교육 슬라이드 생성기
python-pptx로 5개 데크(day1~day5)를 생성. 내용은 placeholder 중심.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 테마 ----------
NAVY   = RGBColor(0x12, 0x20, 0x3C)   # 본문 다크
INK    = RGBColor(0x0F, 0x17, 0x2A)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xF8, 0xFA, 0xFC)
LINE   = RGBColor(0xE2, 0xE8, 0xF0)

# 일자별 강조색 (wayfinding)
DAY_COLORS = {
    1: RGBColor(0x63, 0x66, 0xF1),  # indigo
    2: RGBColor(0x0E, 0xA5, 0xE9),  # sky
    3: RGBColor(0x10, 0xB9, 0x81),  # emerald
    4: RGBColor(0xF5, 0x9E, 0x0B),  # amber
    5: RGBColor(0xF4, 0x3F, 0x5E),  # rose
}
DAY_KR = {1: "월", 2: "화", 3: "수", 4: "목", 5: "금"}

# 3세션/일 시간표 (모든 날 공통): 오전 90분 · 중식 · 오후1 80분 · 휴식 · 오후2 90분
SESSION_MIN   = [90, 80, 90]
SESSION_CLOCK = ["10:00–11:30", "13:00–14:20", "14:30–16:00"]
SESSION_LABEL = ["오전 특강", "오후 특강 1", "오후 특강 2"]

FONT_H = "NanumSquare"        # 제목
FONT_B = "NanumGothic"        # 본문
FONT_M = "NanumGothicCoding"  # 모노/코드

SW, SH = Inches(13.333), Inches(7.5)

# ---------- 저수준 헬퍼 ----------
def _set_font(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp

def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=6, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, font)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold, font) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color
            _set_font(r, font)
    return tb

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    add_rect(slide, 0, 0, SW, SH, fill=color)

def footer(slide, day, total_pages=None):
    # 좌측 푸터 텍스트만 (우측 Day·페이지 표시는 stamp_pages 후처리에서)
    add_text(slide, Inches(0.55), Inches(7.02), Inches(8), Inches(0.4),
             [[("AI 에이전트와 응용  ·  원광대학교 5일 교육", 9, MUTED, False, FONT_B)]],
             anchor=MSO_ANCHOR.MIDDLE)

def stamp_pages(prs, day):
    """빌드 후 모든 슬라이드 우측하단에 'Day N · 요일 · 현재/전체' 표기 (배경색 자동 대응)."""
    total = len(prs.slides._sldIdLst)
    for idx, slide in enumerate(prs.slides, 1):
        dark = False
        try:
            rgb = slide.shapes[0].fill.fore_color.rgb  # 첫 도형 = 배경 사각형
            dark = (rgb[0] + rgb[1] + rgb[2]) < 320
        except Exception:
            pass
        daycol = DAY_COLORS[day]
        pgcol = RGBColor(0xCB, 0xD5, 0xE1) if dark else MUTED
        add_text(slide, Inches(9.83), Inches(7.02), Inches(2.95), Inches(0.4),
                 [[(f"Day {day} · {DAY_KR[day]}", 9, daycol, True, FONT_B),
                   (f"   ·   {idx} / {total}", 9, pgcol, True, FONT_B)]],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# ---------- 슬라이드 타입 ----------
def slide_title(prs, day, day_title, subtitle, sessions, slots):
    s = blank(prs); ac = DAY_COLORS[day]
    bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.35), SH, fill=ac)
    # day chip
    add_rect(s, Inches(0.9), Inches(1.5), Inches(2.5), Inches(0.62), fill=ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.9), Inches(1.5), Inches(2.5), Inches(0.62),
             [[(f"DAY {day}  ·  {DAY_KR[day]}요일", 15, WHITE, True, FONT_H)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(2.45), Inches(11.5), Inches(1.6),
             [[(day_title, 40, WHITE, True, FONT_H)]], anchor=MSO_ANCHOR.TOP, line_spacing=1.05)
    add_text(s, Inches(0.92), Inches(3.95), Inches(11), Inches(0.6),
             [[(subtitle, 17, RGBColor(0xCB,0xD5,0xE1), False, FONT_B)]])
    # session strip
    y = Inches(5.2)
    add_text(s, Inches(0.92), Inches(4.72), Inches(8.5), Inches(0.35),
             [[(f"담당 세션 ({len(sessions)})", 12, ac, True, FONT_B)]])
    n = len(sessions); gap = Inches(0.2)
    cw = Emu(min((int(Inches(11.9)) - (n-1)*int(gap)) // n, int(Inches(3.75))))
    for i, sess in enumerate(sessions):
        sl = slots[i]
        x = Emu(int(Inches(0.9)) + i*(int(cw)+int(gap)))
        add_rect(s, x, y, cw, Inches(1.55), fill=RGBColor(0x1B,0x2B,0x4D),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x+Inches(0.2), y+Inches(0.16), cw-Inches(0.35), Inches(0.35),
                 [[(SESSION_LABEL[sl-1], 13, ac, True, FONT_H)]])
        add_text(s, x+Inches(0.2), y+Inches(0.5), cw-Inches(0.35), Inches(0.3),
                 [[(f"{SESSION_CLOCK[sl-1]} · {SESSION_MIN[sl-1]}분", 10, RGBColor(0x94,0xA3,0xB8), True, FONT_B)]])
        add_text(s, x+Inches(0.2), y+Inches(0.82), cw-Inches(0.35), Inches(0.66),
                 [[(sess, 11, WHITE, False, FONT_B)]], line_spacing=1.0)
    return s

def header(slide, day, kicker, title):
    ac = DAY_COLORS[day]
    bg(slide, WHITE)
    add_rect(slide, 0, 0, Inches(0.18), SH, fill=ac)
    add_text(slide, Inches(0.55), Inches(0.42), Inches(11), Inches(0.35),
             [[(kicker, 12, ac, True, FONT_B)]])
    add_text(slide, Inches(0.55), Inches(0.78), Inches(12.2), Inches(0.8),
             [[(title, 27, INK, True, FONT_H)]])
    add_rect(slide, Inches(0.57), Inches(1.55), Inches(1.1), Inches(0.05), fill=ac)
    footer(slide, day)

def slide_roadmap(prs, day):
    s = blank(prs)
    header(s, day, "PROGRAM ROADMAP", "전체 과정 로드맵")
    titles = {1:"LLM 원리와 이해", 2:"에이전트 원리 + 클로드 코드",
              3:"미니 프로젝트 + 리서치 이론·팀", 4:"팀 리서치 — 기획·구현·실행", 5:"실행·분석 + 보고서·발표"}
    y = Inches(2.4); cw = Inches(2.42); gap = Inches(0.12); x0 = Inches(0.55)
    for d in range(1, 6):
        x = Emu(int(x0) + (d-1)*(int(cw)+int(gap)))
        cur = (d == day)
        col = DAY_COLORS[d]
        add_rect(s, x, y, cw, Inches(2.7), fill=(col if cur else CARD),
                 line=(None if cur else LINE), line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tcol = WHITE if cur else INK
        mcol = RGBColor(0xE2,0xE8,0xF0) if cur else MUTED
        add_text(s, x+Inches(0.2), y+Inches(0.22), cw-Inches(0.4), Inches(0.5),
                 [[(f"DAY {d}", 14, (WHITE if cur else col), True, FONT_H)],
                  [(f"{DAY_KR[d]}요일", 10, mcol, False, FONT_B)]])
        add_text(s, x+Inches(0.2), y+Inches(1.15), cw-Inches(0.4), Inches(1.4),
                 [[(titles[d], 13, tcol, cur, FONT_B)]], line_spacing=1.05)
        if cur:
            add_text(s, x+Inches(0.2), y+Inches(2.25), cw-Inches(0.4), Inches(0.35),
                     [[("● 진행 중", 10, WHITE, True, FONT_B)]])
    add_text(s, Inches(0.55), Inches(5.5), Inches(12), Inches(0.6),
             [[("진행 단계: ", 13, MUTED, True, FONT_B),
               ("이론 → 도구 → 실전 → 결과 → 발표", 13, INK, True, FONT_B)]])
    return s

def slide_agenda(prs, day, sessions, slots):
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, "AGENDA", f"세션 구성 ({len(sessions)})")
    if len(sessions) <= 3:
        y0 = Inches(2.1); rh = Inches(1.4); step = Inches(1.62)
    else:
        y0 = Inches(1.95); rh = Inches(1.15); step = Inches(1.255)
    for i, (t, desc) in enumerate(sessions):
        yy = Emu(int(y0) + i*int(step))
        add_rect(s, Inches(0.55), yy, Inches(12.2), rh, fill=CARD,
                 line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(s, Inches(0.55), yy, Inches(0.14), rh, fill=ac)
        add_rect(s, Inches(0.85), yy+Inches(0.27), Inches(0.62), Inches(0.62), fill=ac,
                 shape=MSO_SHAPE.OVAL)
        add_text(s, Inches(0.85), yy+Inches(0.27), Inches(0.62), Inches(0.62),
                 [[(str(slots[i]), 18, WHITE, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.7), yy+Inches(0.16), Inches(8.9), Inches(0.5),
                 [[(t, 16, INK, True, FONT_H)]])
        add_text(s, Inches(1.7), yy+Inches(0.66), Inches(8.9), Inches(0.42),
                 [[(f"{SESSION_LABEL[slots[i]-1]} · ", 12, ac, True, FONT_B), (desc, 12, MUTED, False, FONT_B)]])
        add_text(s, Inches(10.75), yy+Inches(0.26), Inches(1.85), Inches(0.4),
                 [[(SESSION_CLOCK[slots[i]-1], 12, INK, True, FONT_B)]], align=PP_ALIGN.RIGHT)
        add_text(s, Inches(10.75), yy+Inches(0.64), Inches(1.85), Inches(0.4),
                 [[(f"{SESSION_MIN[slots[i]-1]}분", 12, ac, True, FONT_B)]], align=PP_ALIGN.RIGHT)
    return s

def slide_section(prs, day, sess_no, title, oneliner):
    s = blank(prs); ac = DAY_COLORS[day]
    bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.35), SH, fill=ac)
    add_text(s, Inches(0.95), Inches(2.3), Inches(9), Inches(0.6),
             [[(SESSION_LABEL[sess_no-1], 16, ac, True, FONT_H),
               (f"   {SESSION_CLOCK[sess_no-1]} · {SESSION_MIN[sess_no-1]}분",
                13, RGBColor(0x94,0xA3,0xB8), True, FONT_B)]])
    add_text(s, Inches(0.92), Inches(2.95), Inches(11.4), Inches(1.6),
             [[(title, 36, WHITE, True, FONT_H)]], line_spacing=1.05)
    add_text(s, Inches(0.95), Inches(4.6), Inches(11), Inches(0.8),
             [[(oneliner, 16, RGBColor(0xCB,0xD5,0xE1), False, FONT_B)]], line_spacing=1.1)
    add_rect(s, Inches(0.95), Inches(4.4), Inches(1.0), Inches(0.05), fill=ac)
    return s

def slide_bullets(prs, day, kicker, title, bullets, note=None):
    """bullets: list of (lead, sub_or_None)"""
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, kicker, title)
    y = Inches(2.0)
    for (lead, sub) in bullets:
        add_rect(s, Inches(0.6), y+Inches(0.09), Inches(0.16), Inches(0.16), fill=ac,
                 shape=MSO_SHAPE.OVAL)
        paras = [[(lead, 16, INK, True, FONT_B)]]
        add_text(s, Inches(0.95), y, Inches(11.6), Inches(0.5), paras)
        if sub:
            add_text(s, Inches(0.95), y+Inches(0.42), Inches(11.6), Inches(0.5),
                     [[(sub, 13, MUTED, False, FONT_B)]])
            y = Emu(int(y) + int(Inches(0.95)))
        else:
            y = Emu(int(y) + int(Inches(0.62)))
    if note:
        add_rect(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.72), fill=LIGHT,
                 line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, Inches(0.85), Inches(6.05), Inches(11.6), Inches(0.72),
                 [[("NOTE   ", 12, ac, True, FONT_B), (note, 13, INK, False, FONT_B)]],
                 anchor=MSO_ANCHOR.MIDDLE)
    return s

def slide_twocol(prs, day, kicker, title, left_h, left_items, right_h, right_items):
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, kicker, title)
    for ci, (hh, items, accent) in enumerate([(left_h, left_items, ac), (right_h, right_items, NAVY)]):
        x = Inches(0.55 + ci*6.2)
        add_rect(s, x, Inches(2.0), Inches(5.95), Inches(4.4), fill=CARD,
                 line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(s, x, Inches(2.0), Inches(5.95), Inches(0.7), fill=accent,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x+Inches(0.3), Inches(2.0), Inches(5.4), Inches(0.7),
                 [[(hh, 16, WHITE, True, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
        yy = Inches(2.95)
        for it in items:
            add_text(s, x+Inches(0.35), yy, Inches(5.3), Inches(0.6),
                     [[("•  ", 14, accent, True, FONT_B), (it, 14, INK, False, FONT_B)]], line_spacing=1.0)
            yy = Emu(int(yy) + int(Inches(0.62)))
    return s

def slide_lab(prs, day, title, steps, deliverable):
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, "LAB · 실습", title)
    add_rect(s, Inches(0.55), Inches(1.95), Inches(8.0), Inches(4.55), fill=CARD,
             line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.85), Inches(2.15), Inches(7.4), Inches(0.4),
             [[("실습 단계", 13, ac, True, FONT_B)]])
    yy = Inches(2.7)
    for i, st in enumerate(steps):
        add_rect(s, Inches(0.85), yy, Inches(0.45), Inches(0.45), fill=ac, shape=MSO_SHAPE.OVAL)
        add_text(s, Inches(0.85), yy, Inches(0.45), Inches(0.45),
                 [[(str(i+1), 13, WHITE, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.5), yy+Inches(0.02), Inches(6.9), Inches(0.5),
                 [[(st, 14, INK, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE)
        yy = Emu(int(yy) + int(Inches(0.62)))
    # deliverable card
    add_rect(s, Inches(8.75), Inches(1.95), Inches(4.0), Inches(4.55), fill=NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(9.05), Inches(2.2), Inches(3.4), Inches(0.5),
             [[("산출물", 15, ac, True, FONT_H)]])
    add_text(s, Inches(9.05), Inches(2.8), Inches(3.45), Inches(3.4),
             [[(deliverable, 14, WHITE, False, FONT_B)]], line_spacing=1.2)
    return s

def slide_wrap(prs, day, takeaways, next_teaser):
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, "SUMMARY", "핵심 정리 및 다음 과정")
    add_text(s, Inches(0.6), Inches(1.95), Inches(7), Inches(0.4),
             [[("핵심 정리", 15, ac, True, FONT_H)]])
    yy = Inches(2.5)
    for t in takeaways:
        add_rect(s, Inches(0.62), yy+Inches(0.08), Inches(0.16), Inches(0.16), fill=ac, shape=MSO_SHAPE.OVAL)
        add_text(s, Inches(0.95), yy, Inches(7.0), Inches(0.6),
                 [[(t, 15, INK, False, FONT_B)]], line_spacing=1.05)
        yy = Emu(int(yy) + int(Inches(0.7)))
    add_rect(s, Inches(8.3), Inches(1.95), Inches(4.45), Inches(4.4), fill=CARD,
             line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if day < 5:
        add_text(s, Inches(8.6), Inches(2.2), Inches(3.9), Inches(0.5),
                 [[(f"다음 과정 — Day {day+1}", 15, DAY_COLORS[day+1], True, FONT_H)]])
    else:
        add_text(s, Inches(8.6), Inches(2.2), Inches(3.9), Inches(0.5),
                 [[("수료 · GRADUATION", 14, ac, True, FONT_H)]])
    add_text(s, Inches(8.6), Inches(2.85), Inches(3.85), Inches(3.3),
             [[(next_teaser, 14, INK, False, FONT_B)]], line_spacing=1.25)
    return s

def slide_session_outline(prs, day, sess_no, sess_title, items, duration="약 2시간"):
    """세션 시작 시 in-session 목차. items: list of str (sub-topic)."""
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, f"{SESSION_LABEL[sess_no-1]} · 세션 개요", sess_title)
    # 시간 배지
    add_rect(s, Inches(10.95), Inches(0.5), Inches(1.85), Inches(0.5), fill=LIGHT,
             line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(10.95), Inches(0.5), Inches(1.85), Inches(0.5),
             [[(f"{duration} · {len(items)} 파트", 11, MUTED, True, FONT_B)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.6), Inches(1.85), Inches(8), Inches(0.35),
             [[("이 세션에서 다룰 내용", 13, ac, True, FONT_B)]])
    two_col = len(items) > 6
    per_col = (len(items) + 1) // 2 if two_col else len(items)
    col_w = Inches(5.95) if two_col else Inches(12.2)
    row_h = Inches(0.78) if not two_col else Inches(0.74)
    for i, it in enumerate(items):
        col = 1 if (two_col and i >= per_col) else 0
        row = i - per_col if col == 1 else i
        x = Inches(0.55 + col*6.25)
        yy = Emu(int(Inches(2.35)) + row*int(row_h))
        add_rect(s, x, yy, col_w, row_h-Inches(0.12), fill=CARD,
                 line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(s, x+Inches(0.16), yy+Inches(0.13), Inches(0.42), Inches(0.42), fill=ac, shape=MSO_SHAPE.OVAL)
        add_text(s, x+Inches(0.16), yy+Inches(0.13), Inches(0.42), Inches(0.42),
                 [[(str(i+1), 13, WHITE, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x+Inches(0.78), yy, col_w-Inches(0.95), row_h-Inches(0.12),
                 [[(it, 14, INK, True, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE)
    return s

def slide_process(prs, day, kicker, title, steps, note=None):
    """가로 흐름도. steps: list of (label, sub_or_None), 최대 5."""
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, kicker, title)
    n = len(steps)
    gap = Inches(0.35)
    total_w = Inches(12.2)
    bw = Emu((int(total_w) - (n-1)*int(gap)) // n)
    y = Inches(2.9); bh = Inches(1.9)
    for i, (lab, sub) in enumerate(steps):
        x = Emu(int(Inches(0.55)) + i*(int(bw)+int(gap)))
        add_rect(s, x, y, bw, bh, fill=CARD, line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(s, x, y, bw, Inches(0.12), fill=ac)
        add_text(s, x+Inches(0.15), y+Inches(0.32), bw-Inches(0.3), Inches(0.4),
                 [[(f"{i+1}", 13, ac, True, FONT_H)]])
        add_text(s, x+Inches(0.15), y+Inches(0.72), bw-Inches(0.3), Inches(0.6),
                 [[(lab, 15, INK, True, FONT_B)]], line_spacing=1.0)
        if sub:
            add_text(s, x+Inches(0.15), y+Inches(1.28), bw-Inches(0.3), Inches(0.5),
                     [[(sub, 11, MUTED, False, FONT_B)]], line_spacing=1.0)
        if i < n-1:
            ax = Emu(int(x)+int(bw)+int(Inches(0.05)))
            add_text(s, ax, y+Inches(0.55), Inches(0.3), Inches(0.6),
                     [[("→", 22, ac, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if note:
        add_rect(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(0.72), fill=LIGHT,
                 line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(0.72),
                 [[("NOTE   ", 12, ac, True, FONT_B), (note, 13, INK, False, FONT_B)]],
                 anchor=MSO_ANCHOR.MIDDLE)
    return s

def slide_stat(prs, day, kicker, title, stats, note=None):
    """큰 숫자 강조. stats: list of (big, label), 최대 3."""
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, kicker, title)
    n = len(stats); gap = Inches(0.4)
    bw = Emu((int(Inches(12.2)) - (n-1)*int(gap)) // n)
    y = Inches(2.5); bh = Inches(2.6)
    for i, (big, label) in enumerate(stats):
        x = Emu(int(Inches(0.55)) + i*(int(bw)+int(gap)))
        add_rect(s, x, y, bw, bh, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x+Inches(0.2), y+Inches(0.55), bw-Inches(0.4), Inches(1.2),
                 [[(big, 44, ac, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x+Inches(0.2), y+Inches(1.75), bw-Inches(0.4), Inches(0.7),
                 [[(label, 14, RGBColor(0xCB,0xD5,0xE1), False, FONT_B)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)
    note = note or "수치는 placeholder — 강사가 최신 출처로 보강"
    add_text(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.5),
             [[("NOTE   ", 12, ac, True, FONT_B), (note, 13, MUTED, False, FONT_B)]])
    return s

def slide_diagram(prs, day, kicker, title, caption):
    """도식 placeholder 큰 박스."""
    s = blank(prs); ac = DAY_COLORS[day]
    header(s, day, kicker, title)
    add_rect(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(3.7), fill=CARD,
             line=LINE, line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.55), Inches(3.45), Inches(12.2), Inches(0.7),
             [[("［ 도식 / 다이어그램 placeholder ］", 18, MUTED, True, FONT_B)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.6), Inches(5.95), Inches(12.2), Inches(0.6),
             [[("설명: ", 13, ac, True, FONT_B), (caption, 13, INK, False, FONT_B)]])
    return s

# ---------- 데크 콘텐츠 ----------
def build_day(prs, day, cfg):
    sessions = cfg["sessions"]
    slots = [s.get("slot", i+1) for i, s in enumerate(sessions)]
    slide_title(prs, day, cfg["title"], cfg["subtitle"], cfg["sess_short"], slots)
    slide_roadmap(prs, day)
    if len(sessions) > 1:
        slide_agenda(prs, day, cfg["agenda"], slots)
    for i, blk in enumerate(sessions):
        slot = slots[i]
        slide_section(prs, day, slot, blk["title"], blk["oneliner"])
        if blk.get("outline"):
            slide_session_outline(prs, day, slot, blk["title"], blk["outline"],
                                  duration=f'{SESSION_MIN[slot-1]}분')
        for sl in blk["slides"]:
            kind = sl[0]
            if kind == "bullets":
                slide_bullets(prs, day, sl[1], sl[2], sl[3], sl[4] if len(sl) > 4 else None)
            elif kind == "twocol":
                slide_twocol(prs, day, sl[1], sl[2], sl[3], sl[4], sl[5], sl[6])
            elif kind == "lab":
                slide_lab(prs, day, sl[1], sl[2], sl[3])
            elif kind == "process":
                slide_process(prs, day, sl[1], sl[2], sl[3], sl[4] if len(sl) > 4 else None)
            elif kind == "stat":
                slide_stat(prs, day, sl[1], sl[2], sl[3], sl[4] if len(sl) > 4 else None)
            elif kind == "diagram":
                slide_diagram(prs, day, sl[1], sl[2], sl[3])
    slide_wrap(prs, day, cfg["takeaways"], cfg["teaser"])
    stamp_pages(prs, day)

def new_prs():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    return prs

# 콘텐츠는 별도 모듈에서 가져옴
from slides_content import DAYS

def main():
    outdir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "slides")
    os.makedirs(outdir, exist_ok=True)
    for day in range(1, 6):
        prs = new_prs()
        build_day(prs, day, DAYS[day])
        path = os.path.join(outdir, f"day{day}.pptx")
        prs.save(path)
        print(f"saved {path} ({len(prs.slides._sldIdLst)} slides)")

if __name__ == "__main__":
    main()
